[0]
slide = prs.slides.add_slide(title_slide_layout)

# Ajouter un titre
title = slide.shapes.title
title.text = "Sur l'IA"

# Enregistrer la pr�sentation
prs.save("presentation_ia.pptx")

print("Pr�sentation cr��e avec succ�s : presentation_ia.pptx")
```

Pour ex�cuter ce code, assurez-vous d'avoir install� la biblioth�que `python-pptx`.  Vous pouvez l'installer via pip :

```bash
pip install python-pptx
```

Apr�s l'ex�cution, un fichier nomm� `presentation_ia.pptx` sera cr�� dans le m�me r�pertoire que votre script Python.  Ce fichier contiendra une pr�sentation PowerPoint avec une diapositive ayant comme titre "Sur l'IA".


Pour ajouter plus de contenu (sous-titres, images, puces, etc.),  vous devrez explorer davantage les fonctionnalit�s de la biblioth�que `python-pptx`.  La documentation officielle est une excellente ressource :  [https://python-pptx.readthedocs.io/en/latest/](https://python-pptx.readthedocs.io/en/latest/)


Voici un exemple plus complet ajoutant un sous-titre et un paragraphe de texte:

```python
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)

title = slide.shapes.title
title.text = "Sur l'IA"

subtitle = slide.placeholders[1] # Placeholder pour le sous-titre
subtitle.text = "Introduction � l'intelligence artificielle"

# Ajouter un paragraphe de texte
left = Inches(1)
top = Inches(2.5)
width = Inches(6)
height = Inches(2)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.text = "L'intelligence artificielle (IA) est un domaine vaste et complexe..."
p = tf.paragraphs[0]